from io import BytesIO
from docx import Document
from pptx import Presentation
from typing import Union, IO
from pypdf import PdfReader
import pandas as pd
import tiktoken
import os
import datetime
import json

class File:
    def __init__(self) -> None:
        pass
    
    def tokenize(self,text):
        enc = tiktoken.encoding_for_model("gpt-4o")
        tokens = enc.encode(text)

        return tokens
    
    def load_document(self, file: Union[str, IO[bytes]] = None):
        #TODO: add more file types
        #I guess it supports doc and the varities of .pp**/.xl**, but needs to test it
        if isinstance(file, str):
            ext = os.path.splitext(file)[1].lower()
        else:
            ext = os.path.splitext(file.name)[1].lower()

        if ext == ".docx":
            return self.load_word(file)
        elif ext == ".xlsx":
            return self.load_excel(file)
        elif ext == ".pptx":
            return self.load_pptx(file)
        elif ext == ".pdf":
            return self.load_pdf(file)
        else:
            raise ValueError("Tipo de arquivo não suportado. Por favor, insira um arquivo .docx, .xlsx ou .pptx.")

    def load_excel(self, file: str | IO[bytes ]= None):
        pass # too lazy to do this. Only when i need it
    
    def load_pptx(self, file: str | IO[bytes ]= None):
        try:
            prs = Presentation(file)
        except Exception as e:
            raise(e)

        properties_dir = dir(prs.core_properties)

        properties_keys = [k for k in properties_dir if not k.startswith("_")]
        
        properties = {}
        for key in properties_keys:
            properties[key] = prs.core_properties.__getattribute__(key)

        properties["slide count"] = len(prs.slides)

        tokens = []
        final_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if isinstance(shape.text, (int,float)) or not shape.text.strip() or shape.text.isdigit():
                        continue
                    final_text += shape.text
                    pretokens = self.tokenize(shape.text)
                    for token in pretokens:
                        tokens.append(token)

                if shape.has_table:
                    table = shape.table 
                    table_data = []
                    for row in table.rows:
                        for cell in row.cells:
                            if isinstance(cell, (int,float)) or not cell.text.strip():
                                    continue
                            final_text += cell.text.rstrip()
                            pretokens = self.tokenize(cell.text.rstrip())
                            for token in pretokens:
                                tokens.append(token)
                if shape.has_chart:
                    try:
                        blob_stream = BytesIO(shape.chart._workbook.xlsx_part.blob)
                        df = pd.read_excel(blob_stream)
                        blob_stream.close()
                        
                        columns = df.columns.tolist()
                        for index, column in enumerate(columns):
                            if not column.strip() or column.isdigit() or isinstance(column, (int,float)):
                                    continue
                            final_text += column
                            
                            pretokens = self.tokenize(column)
                            for token in pretokens:
                                tokens.append(token)
                            
                        row_data = df.values.tolist()
                        for row in row_data:
                            for index, data in enumerate(row):
                                if not data.strip() or data.isdigit() or isinstance(data, (int,float)):
                                    continue
                                
                                final_text += data
                            
                                pretokens = self.tokenize(column)
                                for token in pretokens:
                                    tokens.append(token)
                        
                        row_data = None
                        columns = None
                        df = None
                        
                    except Exception as e:
                        continue
                                
        properties["word count"] = len(final_text.split(" "))

        # properties["tokens"] = tokens
        
        try:
            properties["created"] = datetime.datetime.isoformat(properties["created"])
            properties["last_printed"] = datetime.datetime.isoformat(properties["last_printed"])
            properties["modified"] = datetime.datetime.isoformat(properties["modified"])
        except:
            pass
        
        properties["tokens count"] = len(tokens)

        result = json.dumps(obj=properties,skipkeys=True, default=lambda o: '<not serializable>',indent=2,ensure_ascii=False)
        
        # print(json.loads(result))

        # The reason behind it returning the dict loaded from the json and not the json object itself
        # is that in the rust side it is programmed to handle dict files
        # was too lazy to fix that so i am returning the dict after parsing it to json
        return json.loads(result)

    def load_word(self, file: str | IO[bytes ]= None):
        try:
            file = Document(file)
        except Exception as e:
            raise(e)

        properties_dir = dir(file.core_properties)
        
        properties_keys = [k for k in properties_dir if not k.startswith("_")]
        
        properties = {}
        for key in properties_keys:
            properties[key] = file.core_properties.__getattribute__(key)

        paragraphs_text = ""
        tokens = []
        for paragraph in file.paragraphs:
            text = paragraph.text.rstrip()
            text = text.lstrip()
            paragraphs_text += " "+text
            pretokens = self.tokenize(paragraph.text)
            for token in pretokens:
                tokens.append(token)

        words = paragraphs_text.split(" ")

        while "" in words:
            words.remove("")

        properties["word_count"] = len(words)

        # properties["tokens"] = tokens
        
        
        try:
            properties["created"] = datetime.datetime.isoformat(properties["created"])
            properties["last_printed"] = datetime.datetime.isoformat(properties["last_printed"])
            properties["modified"] = datetime.datetime.isoformat(properties["modified"])
        except:
            pass
        properties["tokens_count"] = len(tokens)
        
        result = json.dumps(obj=properties,skipkeys=True, default=lambda o: '<not serializable>',indent=2,ensure_ascii=False)
        
        return json.loads(result)
    
    def load_pdf(self, file: str | IO[bytes ]= None, pdf_password: str = None):
        reader = PdfReader(file)
        number_of_pages = reader.get_num_pages()
        
        metadata = reader.metadata

        text = ""

        for page in reader.pages:
            text += page.extract_text()

        properties = {}
        
        properties["Creator"] = metadata.creator
        properties["Creation Date"] = datetime.datetime.isoformat(metadata.creation_date)
        properties["Modification Date"] = datetime.datetime.isoformat(metadata.modification_date)
        
        words = text.split(" ")

        while "" in words:
            words.remove("")
        
        tokens = self.tokenize(text)
        
        properties["number of pages"] = number_of_pages
        properties["word count"] = len(words)
        
        # properties["tokens"] = tokens
        properties["token count"] = len(tokens)
                
        return properties
